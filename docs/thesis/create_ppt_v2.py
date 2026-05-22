import sys
sys.path.insert(0, r'D:\SwapU\docs\thesis\pip_libs')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
def OxmlElement(tag, **kwargs):
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
             'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
             'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
    prefix, local = tag.split(':') if ':' in tag else ('a', tag)
    ns = nsmap.get(prefix, nsmap['a'])
    return etree.SubElement(etree.Element('dummy'), f'{{{ns}}}{local}')
import os, copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

C1 = RGBColor(0x0F, 0x17, 0x2A)
C2 = RGBColor(0x1B, 0x2A, 0x4A)
C3 = RGBColor(0x0D, 0x94, 0x88)
C4 = RGBColor(0x14, 0xB8, 0xA6)
C5 = RGBColor(0x5E, 0xEA, 0xD4)
C6 = RGBColor(0xF0, 0xFD, 0xFA)
C7 = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LGRAY = RGBColor(0x94, 0xA3, 0xB8)
RED = RGBColor(0xEF, 0x44, 0x44)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xF9, 0x73, 0x16)

IMG = r'd:\SwapU\docs\thesis\thesis_images'

def bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill, line=None, radius=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(1)
    else: s.line.fill.background()
    if radius:
        sp = s._element.spPr
        prstGeom = sp.find(qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = OxmlElement('a:avLst')
                prstGeom.append(avLst)
            gd = OxmlElement('a:gd')
            gd.set('name', 'adj')
            gd.set('fmla', f'val {radius}')
            avLst.append(gd)
    return s

def txt(slide, x, y, w, h, text, sz=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font='Calibri', spacing=None, valign=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if valign: tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold; p.font.name = font
    p.alignment = align
    if spacing: p.space_after = Pt(spacing)
    return tb

def rich_txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, valign=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    if valign: tf.vertical_anchor = valign
    for i, (text, sz, color, bold, font) in enumerate(runs):
        if i == 0: p = tf.paragraphs[0]
        else: p = tf.add_paragraph()
        p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.bold = bold; p.font.name = font; p.alignment = align
        p.space_after = Pt(4)
    return tb

def bullets(slide, x, y, w, h, items, sz=14, color=DARK, font='Calibri', spacing=6, bullet_color=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0: p = tf.paragraphs[0]
        else: p = tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]; p.font.size = Pt(item[1]); p.font.color.rgb = item[2]; p.font.bold = item[3] if len(item)>3 else False
        else:
            p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = font; p.space_after = Pt(spacing)
        pPr = p._p.get_or_add_pPr()
        buNone = pPr.find(qn('a:buNone'))
        if buNone is not None: pPr.remove(buNone)
        buChar = OxmlElement('a:buChar'); buChar.set('char', '\u25B8')
        pPr.append(buChar)
        if bullet_color:
            buClr = OxmlElement('a:buClr')
            srgb = OxmlElement('a:srgbClr'); srgb.set('val', str(bullet_color))
            buClr.append(srgb); pPr.append(buClr)
    return tb

def img(slide, name, x, y, w=None, h=None):
    path = os.path.join(IMG, name)
    if not os.path.exists(path): return
    kw = {}
    if w: kw['width'] = Inches(w)
    if h: kw['height'] = Inches(h)
    slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)

def accent_bar(slide, x, y, w=0.06, h=0.5):
    rect(slide, x, y, w, h, C3)

def section_header(slide, num, title, subtitle=''):
    rect(slide, 0, 0, 13.333, 0.06, C3)
    txt(slide, 0.8, 0.35, 1.0, 0.5, num, sz=36, color=C3, bold=True, font='Arial Black')
    txt(slide, 1.8, 0.4, 8, 0.5, title, sz=28, color=C1, bold=True)
    if subtitle:
        txt(slide, 1.8, 0.95, 8, 0.3, subtitle, sz=13, color=GRAY)
    rect(slide, 0.8, 1.25, 2.0, 0.03, C3)

def card(slide, x, y, w, h, accent=True):
    rect(slide, x, y, w, h, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
    if accent:
        rect(slide, x+0.15, y+0.2, 0.05, h-0.4, C3)

def page_num(slide, n, total):
    txt(slide, 12.0, 7.1, 1.0, 0.3, f'{n}/{total}', sz=10, color=LGRAY, align=PP_ALIGN.RIGHT)

TOTAL = 22

# ════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C1)
rect(s, 0, 0, 13.333, 0.06, C3)
rect(s, 0, 7.44, 13.333, 0.06, C3)
rect(s, 1.0, 2.0, 0.06, 2.5, C3)
txt(s, 1.4, 2.0, 10, 0.8, '基于RAG的校园智能二手交易系统', sz=40, color=WHITE, bold=True, font='Calibri')
txt(s, 1.4, 2.85, 10, 0.5, '设计与实现', sz=28, color=C4, font='Calibri')
txt(s, 1.4, 3.6, 10, 0.4, 'SwapU \u2014 让校园闲置流转更智能', sz=16, color=C5)
rect(s, 1.4, 4.3, 3.0, 0.02, C3)
txt(s, 1.4, 4.6, 6, 0.35, '答辩人：刘遥杰', sz=16, color=C4)
txt(s, 1.4, 5.0, 6, 0.35, '指导教师：霞', sz=16, color=C4)
txt(s, 1.4, 5.5, 8, 0.3, '武汉理工大学  计算机与人工智能学院  计算机科学与技术 zy2201班', sz=13, color=LGRAY)
txt(s, 9.0, 6.5, 3.5, 0.3, '2026年5月', sz=13, color=LGRAY, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 2: OUTLINE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
rect(s, 0, 0, 13.333, 0.06, C3)
txt(s, 0.8, 0.4, 5, 0.5, '目  录', sz=32, color=C1, bold=True)
rect(s, 0.8, 0.95, 1.5, 0.03, C3)
secs = [
    ('01','选题背景与意义','校园闲置错配痛点\n现有平台三方面错位','img_d2d08876.png'),
    ('02','系统架构设计','五层架构 + 四种存储引擎\n技术栈选型依据','img_ce8dc44d.png'),
    ('03','核心功能展示','AI润色 / RAG导购\n订单交易 / 即时通讯','img_74beebc6.png'),
    ('04','关键技术实现','JWT / AES / 双轨同步\n订单状态机 / ES搜索','img_dd555ab5.png'),
    ('05','系统测试','功能 / 性能 / 安全\n三维度验证','img_8370a17d.png'),
    ('06','总结与展望','主要工作与创新点\n改进方向','img_a1392b1a.png'),
]
for i,(num,title,desc,_) in enumerate(secs):
    r = i // 3; c = i % 3
    x = 0.8 + c*4.1; y = 1.5 + r*2.7
    rect(s, x, y, 3.8, 2.4, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=8000)
    rect(s, x, y, 3.8, 0.5, C2, radius=8000)
    txt(s, x+0.2, y+0.05, 0.6, 0.4, num, sz=22, color=C3, bold=True, font='Arial Black')
    txt(s, x+0.8, y+0.08, 2.8, 0.35, title, sz=16, color=WHITE, bold=True)
    txt(s, x+0.2, y+0.65, 3.4, 1.5, desc, sz=12, color=GRAY)
page_num(s, 2, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 3: BACKGROUND - PROBLEM
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '01', '选题背景与意义', '为什么需要 SwapU？')
txt(s, 0.8, 1.5, 5.5, 0.4, '现实痛点', sz=20, color=C3, bold=True)
rect(s, 0.8, 1.95, 5.5, 2.5, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 2.1, 0.05, 2.2)
bullets(s, 1.3, 2.1, 4.8, 2.2, [
    '每学期末大量教材、电子设备被淘汰',
    '新生群体年复一年存在刚性采购需求',
    '二者之间存在结构性错配',
    '构建校园二手交易生态具有现实意义',
], sz=13, bullet_color=C3)

txt(s, 0.8, 4.7, 5.5, 0.4, '现有平台三方面错位', sz=20, color=C3, bold=True)
rect(s, 0.8, 5.15, 5.5, 2.0, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 5.3, 0.05, 1.7)
bullets(s, 1.3, 5.3, 4.8, 1.7, [
    '信息匹配效率不足 — 关心相邻宿舍楼而非全国信用分',
    '交易摩擦过大 — 校园当面交易为主，平台缺乏面交支持',
    '发布门槛偏高 — 学生不愿为短期闲置品撰写详细描述',
], sz=13, bullet_color=C3)

txt(s, 7.0, 1.5, 5.5, 0.4, 'SwapU 的解决思路', sz=20, color=C3, bold=True)
solutions = [
    ('发布侧','AI 一键润色','DeepSeek Chat 生成结构化商品描述\nSystem Prompt 十余版本迭代，合格率 40% \u2192 95%'),
    ('导购侧','RAG 智能导购','先从 PgVector 检索真实在售商品\n再注入上下文交由模型推荐，从源头防幻觉'),
    ('交易侧','校园当面交易','WebSocket 即时沟通\n支付宝沙箱完成支付闭环'),
]
for i,(side,title,desc) in enumerate(solutions):
    y = 2.0 + i*1.7
    rect(s, 7.0, y, 5.5, 1.5, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
    rect(s, 7.0, y, 1.2, 1.5, C2, radius=5000)
    txt(s, 7.1, y+0.15, 1.0, 0.3, side, sz=11, color=C4, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 7.1, y+0.5, 1.0, 0.3, title, sz=11, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, 8.4, y+0.15, 3.9, 1.2, desc, sz=12, color=GRAY)
page_num(s, 3, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 4: RESEARCH STATUS
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '01', '国内外研究现状', '从学术验证到工程落地的跨越')
txt(s, 0.8, 1.5, 5.8, 0.4, '二手交易平台演进', sz=18, color=C3, bold=True)
rect(s, 0.8, 1.95, 5.8, 2.2, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 2.1, 0.05, 1.9)
bullets(s, 1.3, 2.1, 5.2, 1.9, [
    '国外：eBay 信誉评分 + Craigslist 轻量发布 + Poshmark 社交电商',
    '国内：闲鱼芝麻信用体系 + 转转3C质检验机',
    '共同范式：C2C 大市场撮合，缺乏校园场景针对性设计',
], sz=12, bullet_color=C3)

txt(s, 0.8, 4.4, 5.8, 0.4, '搜索技术演进', sz=18, color=C3, bold=True)
rect(s, 0.8, 4.85, 5.8, 2.2, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 5.0, 0.05, 1.9)
bullets(s, 1.3, 5.0, 5.2, 1.9, [
    'MySQL LIKE \u2192 Elasticsearch BM25 + 协同过滤',
    '关键词匹配无法理解口语化查询（如"考研资料"）',
    '大语言模型突破语义理解困境，但存在幻觉问题',
], sz=12, bullet_color=C3)

txt(s, 7.0, 1.5, 5.5, 0.4, 'RAG 技术原理', sz=18, color=C3, bold=True)
rect(s, 7.0, 1.95, 5.5, 2.2, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 7.2, 2.1, 0.05, 1.9)
bullets(s, 7.5, 2.1, 4.8, 1.9, [
    '核心思想：先检索后生成 [Lewis et al., 2020]',
    '从外部知识库检索最相关真实数据',
    '注入 Prompt 约束模型输出在真实数据范围内',
    '从根本上缓解大模型幻觉问题',
], sz=12, bullet_color=C3)

txt(s, 7.0, 4.4, 5.5, 0.4, 'SwapU 的定位', sz=18, color=C3, bold=True)
rect(s, 7.0, 4.85, 5.5, 2.2, C2, radius=5000)
txt(s, 7.3, 5.1, 5.0, 1.8,
    '不满足于实验室验证 RAG 可行性\n'
    '而是部署到真实校园交易场景\n'
    '面对商品库持续变动、用户查询口语化\n'
    '模型输出不可控等工程挑战\n'
    '逐一给出解决方案',
    sz=14, color=WHITE)
page_num(s, 4, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 5: ARCHITECTURE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '02', '系统架构设计', '前后端分离 B/S 架构，五层结构')
img(s, 'img_ce8dc44d.png', 0.3, 1.4, h=5.2)
txt(s, 7.2, 1.4, 5.5, 0.4, '五层架构详解', sz=18, color=C3, bold=True)
layers = [
    ('表现层','Vue 3 + Element Plus + Pinia\nNginx 反向代理 /api \u2192 后端8080','C3'),
    ('网关控制层','LoginInterceptor JWT鉴权\nAdminInterceptor role=1校验\nGlobalLogAspect AOP审计','C4'),
    ('业务逻辑层','15 Controller + 13 ServiceImpl\nItemServiceImpl 依赖最密集\nTradeOrderServiceImpl 事务密度最高','C5'),
    ('数据持久层','MySQL ACID写入 | Redis 缓存\nES 全文检索 | PgVector 向量检索','C3'),
    ('基础设施层','RocketMQ 消息驱动 | 支付宝沙箱\nOllama bge-m3 | DeepSeek Chat','C4'),
]
for i,(name,desc,clr) in enumerate(layers):
    y = 1.9 + i*1.05
    rect(s, 7.2, y, 5.5, 0.9, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=3000)
    rect(s, 7.2, y, 0.06, 0.9, C3)
    txt(s, 7.5, y+0.02, 1.8, 0.3, name, sz=13, color=C3, bold=True)
    txt(s, 7.5, y+0.32, 5.0, 0.55, desc, sz=10, color=GRAY)
page_num(s, 5, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 6: TECH STACK
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '02', '技术栈选型', '各技术选型依据与职责')
techs = [
    ('后端框架','Spring Boot 3.2','自动配置+起步依赖\n降低项目搭建成本','C3'),
    ('前端框架','Vue 3 + Vite','Composition API\nPinia 状态管理','C4'),
    ('ORM框架','MyBatis-Plus','BaseMapper + IService\n简化数据库操作','C5'),
    ('搜索引擎','Elasticsearch 8.12','IK分词 + BoolQuery\n高亮标注 + 自动补全','C3'),
    ('向量数据库','PgVector','bge-m3 1024维向量\n余弦相似度检索','C4'),
    ('消息队列','RocketMQ 5.1','延时消息30分钟超时\n异步同步ES+PgVector','C5'),
    ('AI框架','Spring AI','ChatClient抽象层\n模型切换无需改业务代码','C3'),
    ('大语言模型','DeepSeek Chat','文案润色 + 智能导购\ntemperature=0.1润色','C4'),
    ('支付','支付宝沙箱','RSA验签 + 幂等性检查\n异步回调更新状态','C5'),
    ('缓存','Redis 7.0','热点商品详情缓存\n命中率 > 95%','C3'),
    ('数据库','MySQL 8.0','ACID事务 + 关联查询\n系统主数据源','C4'),
    ('Embedding','Ollama bge-m3','本地部署无需GPU\n1024维向量生成','C5'),
]
for i,(cat,name,desc,clr) in enumerate(techs):
    r = i // 4; c = i % 4
    x = 0.6 + c*3.15; y = 1.5 + r*1.9
    rect(s, x, y, 2.95, 1.7, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
    rect(s, x, y, 2.95, 0.4, C3, radius=5000)
    txt(s, x+0.15, y+0.05, 2.6, 0.3, f'{cat}  {name}', sz=11, color=WHITE, bold=True)
    txt(s, x+0.15, y+0.5, 2.6, 1.1, desc, sz=10, color=GRAY)
page_num(s, 6, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 7: UI - Login & Home
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '03', '核心功能展示 \u2014 登录与首页', '用户认证 + 商品浏览')
img(s, 'img_f93bbfac.png', 0.5, 1.5, h=4.8)
img(s, 'img_a1b8e6be.png', 6.8, 1.5, h=4.8)
rect(s, 0.5, 6.5, 5.8, 0.7, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=3000)
accent_bar(s, 0.7, 6.6, 0.04, 0.5)
txt(s, 1.0, 6.55, 5.0, 0.6, 'BCrypt密码加密 + JWT Token签发\n普通用户/管理员双角色登录', sz=11, color=GRAY)
rect(s, 6.8, 6.5, 5.8, 0.7, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=3000)
accent_bar(s, 7.0, 6.6, 0.04, 0.5)
txt(s, 7.3, 6.55, 5.0, 0.6, '卡片式信息流 + 分类导航 + ES搜索\n最新/价格/热度四种排序', sz=11, color=GRAY)
page_num(s, 7, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: UI - AI Features
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '03', '核心功能展示 \u2014 AI 智能交互', '发布侧润色 + 导购侧RAG')
img(s, 'img_6f39e11f.png', 0.3, 1.5, h=4.5)
img(s, 'img_74beebc6.png', 6.5, 1.5, h=4.5)
rect(s, 0.3, 6.2, 6.0, 0.9, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=3000)
accent_bar(s, 0.5, 6.3, 0.04, 0.7)
txt(s, 0.8, 6.25, 5.3, 0.8, 'AI文案润色 /ai/polish\n输入关键词 \u2192 DeepSeek生成结构化描述\ntemperature=0.1，合格率95%+', sz=11, color=GRAY)
rect(s, 6.5, 6.2, 6.5, 0.9, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=3000)
accent_bar(s, 6.7, 6.3, 0.04, 0.7)
txt(s, 7.0, 6.25, 5.8, 0.8, 'AI智能导购 /ai/chat\n自然语言 \u2192 PgVector检索topK=3 \u2192 RAG推荐\n返回AI文本 + 可点击商品卡片', sz=11, color=GRAY)
page_num(s, 8, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: UI - Chat & Orders
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '03', '核心功能展示 \u2014 即时通讯与订单', 'WebSocket + 支付宝沙箱')
img(s, 'img_ae0637aa.png', 0.3, 1.5, h=4.5)
img(s, 'img_a0a7b452.png', 6.5, 1.5, h=4.5)
rect(s, 0.3, 6.2, 6.0, 0.9, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=3000)
accent_bar(s, 0.5, 6.3, 0.04, 0.7)
txt(s, 0.8, 6.25, 5.3, 0.8, 'WebSocket即时通讯\n文本+图片消息 | 未读红点标识\n消息持久化至数据库', sz=11, color=GRAY)
rect(s, 6.5, 6.2, 6.5, 0.9, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=3000)
accent_bar(s, 6.7, 6.3, 0.04, 0.7)
txt(s, 7.0, 6.25, 5.8, 0.8, '订单管理\n买入/卖出双视角 | 确认收货/评价\nRocketMQ 30分钟超时自动取消', sz=11, color=GRAY)
page_num(s, 9, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: UI - Admin & ER
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '03', '核心功能展示 \u2014 后台管理与数据库', 'ECharts + 五表ER设计')
img(s, 'img_ecb8add6.png', 0.3, 1.5, h=4.5)
img(s, 'img_b4942988.png', 6.5, 1.5, h=4.5)
rect(s, 0.3, 6.2, 6.0, 0.9, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=3000)
accent_bar(s, 0.5, 6.3, 0.04, 0.7)
txt(s, 0.8, 6.25, 5.3, 0.8, '管理员Dashboard\nECharts运营指标 + 商品审核 + 举报处理\n举报处理联动更新物品状态+系统通知', sz=11, color=GRAY)
rect(s, 6.5, 6.2, 6.5, 0.9, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=3000)
accent_bar(s, 6.7, 6.3, 0.04, 0.7)
txt(s, 7.0, 6.25, 5.8, 0.8, '核心数据库ER图\n5个核心表：sys_user / item / trade_order / comment / report\n第三范式 + 受控冗余（order冗余amount快照）', sz=11, color=GRAY)
page_num(s, 10, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 11: KEY TECH - JWT
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '04', '关键技术 \u2014 JWT无状态认证', '为什么选JWT？完整认证流程')
rect(s, 0.8, 1.5, 5.8, 2.5, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 1.65, 0.05, 2.2)
txt(s, 1.3, 1.55, 5.0, 0.35, '选型理由', sz=16, color=C3, bold=True)
bullets(s, 1.3, 1.95, 5.0, 2.0, [
    'JWT无状态，水平扩展无需共享Session',
    'Session/Cookie需额外Redis Session，增加复杂度',
    'Token自包含用户信息，服务端无需查询会话存储',
    '任何实例均可独立验证Token有效性',
], sz=12, bullet_color=C3)

rect(s, 0.8, 4.2, 5.8, 3.0, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 4.35, 0.05, 2.7)
txt(s, 1.3, 4.25, 5.0, 0.35, '认证流程', sz=16, color=C3, bold=True)
bullets(s, 1.3, 4.65, 5.0, 2.5, [
    '1. 用户登录 \u2192 JwtUtils.generateToken() 签发',
    '2. 前端localStorage存储Token',
    '3. Axios拦截器注入 Authorization 请求头',
    '4. LoginInterceptor.validateToken() 验证',
    '5. 解析Claims \u2192 userId/username/role 存入Request域',
    '6. Integer/Long类型兼容（JVM版本差异）',
], sz=12, bullet_color=C3)

rect(s, 7.0, 1.5, 5.5, 5.7, C2, radius=5000)
txt(s, 7.3, 1.7, 5.0, 0.35, '关键代码', sz=16, color=C3, bold=True)
code_text = (
    '// JwtUtils.java \u2014 HMAC-SHA256\u7b7e\u540d\n'
    'SecretKey key = Keys.hmacShaKeyFor(\n'
    '    SECRET.getBytes(UTF_8));\n'
    '// Token\u6709\u6548\u671f24\u5c0f\u65f6\n'
    'long EXPIRE = 24*60*60*1000L;\n\n'
    '// LoginInterceptor.java\n'
    'String token = request\n'
    '    .getHeader("Authorization");\n'
    'Claims claims =\n'
    '    JwtUtils.validateToken(token);\n'
    '// Integer\u2192Long\u7c7b\u578b\u517c\u5bb9\n'
    'if (userIdObj instanceof Integer) {\n'
    '  request.setAttribute("userId",\n'
    '    ((Integer)userIdObj).longValue());\n'
    '}'
)
txt(s, 7.3, 2.2, 5.0, 4.5, code_text, sz=11, color=C5, font='Consolas')
page_num(s, 11, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 12: KEY TECH - AES
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '04', '关键技术 \u2014 AES字段加密', 'TypeHandler机制，业务零感知')
rect(s, 0.8, 1.5, 5.8, 2.8, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 1.65, 0.05, 2.5)
txt(s, 1.3, 1.55, 5.0, 0.35, '加密方案', sz=16, color=C3, bold=True)
bullets(s, 1.3, 1.95, 5.0, 2.2, [
    'MyBatis-Plus TypeHandler 机制',
    '写入时自动加密：AesEncryptUtils.encrypt()',
    '读取时自动解密：AesEncryptUtils.decrypt()',
    '注解：@TableField(typeHandler = AesEncryptTypeHandler.class)',
    '业务代码无需感知加解密过程',
], sz=12, bullet_color=C3)

rect(s, 0.8, 4.5, 5.8, 2.5, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 4.65, 0.05, 2.2)
txt(s, 1.3, 4.55, 5.0, 0.35, '安全效果', sz=16, color=C3, bold=True)
bullets(s, 1.3, 4.95, 5.0, 2.0, [
    '数据库 phone 字段存储为 Base64 密文',
    'API 返回时自动解密为明文',
    '即使数据库泄露也无法获取手机号',
    'AES-128 对称加密，密钥配置在 application.yml',
], sz=12, bullet_color=C3)

rect(s, 7.0, 1.5, 5.5, 5.5, C2, radius=5000)
txt(s, 7.3, 1.7, 5.0, 0.35, '关键代码', sz=16, color=C3, bold=True)
code_text = (
    '// AesEncryptTypeHandler.java\n'
    'public void setNonNullParameter(\n'
    '    PreparedStatement ps, int i,\n'
    '    String param, JdbcType type) {\n'
    '  // \u5199\u5165\u65f6\u52a0\u5bc6\n'
    '  ps.setString(i,\n'
    '    AesEncryptUtils.encrypt(param));\n'
    '}\n\n'
    'public String getNullableResult(\n'
    '    ResultSet rs, String col) {\n'
    '  String val = rs.getString(col);\n'
    '  // \u8bfb\u53d6\u65f6\u89e3\u5bc6\n'
    '  return AesEncryptUtils\n'
    '    .decrypt(val);\n'
    '}\n\n'
    '// Entity\u6ce8\u89e3\n'
    '@TableField(typeHandler =\n'
    '  AesEncryptTypeHandler.class)\n'
    'private String phone;'
)
txt(s, 7.3, 2.2, 5.0, 4.5, code_text, sz=11, color=C5, font='Consolas')
page_num(s, 12, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 13: KEY TECH - RAG
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '04', '关键技术 \u2014 RAG智能导购', '三步链路：检索 \u2192 拼接 \u2192 推理')
img(s, 'img_a1392b1a.png', 0.3, 1.4, h=3.5)
steps = [
    ('1','向量检索','bge-m3 生成1024维向量 \u2192 PgVector余弦相似度\ntopK=3, similarityThreshold=0.0\n校园商品库有限，阈值0.0避免查询无果'),
    ('2','上下文拼接','检索结果注入 System Prompt\n历史对话截取最近6条（3轮）\n检索为空时不强行推荐'),
    ('3','推理与封装','chatClient.call() 同步推理\n返回 AI文本 + 商品卡片数据\n平均耗时3.5s（检索200ms+推理3s）'),
]
for i,(num,title,desc) in enumerate(steps):
    y = 1.5 + i*1.85
    rect(s, 7.0, y, 5.8, 1.65, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
    rect(s, 7.0, y, 0.7, 1.65, C2, radius=5000)
    txt(s, 7.05, y+0.2, 0.6, 0.5, num, sz=28, color=C3, bold=True, font='Arial Black', align=PP_ALIGN.CENTER)
    txt(s, 7.8, y+0.1, 4.8, 0.35, title, sz=16, color=C3, bold=True)
    txt(s, 7.8, y+0.5, 4.8, 1.0, desc, sz=11, color=GRAY)

txt(s, 7.0, 7.0, 5.8, 0.35, '防幻觉：Prompt硬性指令 + RAG事实锚点', sz=13, color=C3, bold=True, align=PP_ALIGN.CENTER)
page_num(s, 13, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 14: KEY TECH - DUAL SYNC
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '04', '关键技术 \u2014 RAG双轨同步策略', '增量保实时，全量保正确')
img(s, 'img_dd555ab5.png', 0.3, 1.4, h=3.8)
rect(s, 7.0, 1.5, 5.8, 2.5, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
rect(s, 7.0, 1.5, 0.7, 2.5, C3, radius=5000)
txt(s, 7.05, 1.6, 0.6, 0.4, '\u2160', sz=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 7.9, 1.6, 4.7, 0.35, '第一轨：增量同步', sz=16, color=C3, bold=True)
bullets(s, 7.9, 2.0, 4.7, 1.8, [
    'RocketMQ 消息驱动',
    'MySQL事务完成 \u2192 发送ItemSyncMessage',
    'ItemSyncListener消费 \u2192 ES+PgVector同步',
    'MQ不可用时降级为同步直连',
], sz=11, bullet_color=C3)

rect(s, 7.0, 4.2, 5.8, 2.5, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
rect(s, 7.0, 4.2, 0.7, 2.5, C4, radius=5000)
txt(s, 7.05, 4.3, 0.6, 0.4, '\u2161', sz=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 7.9, 4.3, 4.7, 0.35, '第二轨：全量同步', sz=16, color=C3, bold=True)
bullets(s, 7.9, 4.7, 4.7, 1.8, [
    '管理员触发 /admin/es/sync',
    '遍历所有 status=1 的商品',
    '批量写入 ES 和 PgVector',
    '用于系统初始化或灾难恢复',
], sz=11, bullet_color=C3)

rect(s, 0.8, 5.5, 5.8, 1.5, C2, radius=5000)
txt(s, 1.1, 5.7, 5.3, 1.0,
    '核心思想\n增量保实时 \u2014 日常运行时数据时效性\n全量保正确 \u2014 兜底修复可能的不一致',
    sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
page_num(s, 14, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 15: KEY TECH - ORDER STATE MACHINE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '04', '关键技术 \u2014 订单状态机与支付', '编程式事务 + 支付安全三步验证')
img(s, 'img_8370a17d.png', 0.3, 1.4, h=3.0)
img(s, 'img_2514cc9f.png', 6.5, 1.4, h=3.0)

rect(s, 0.8, 4.6, 5.8, 2.6, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 4.75, 0.05, 2.3)
txt(s, 1.3, 4.65, 5.0, 0.35, '编程式事务 TransactionTemplate', sz=15, color=C3, bold=True)
bullets(s, 1.3, 5.05, 5.0, 2.0, [
    '业务校验通过后才开启事务，避免长事务',
    '@Transactional 是AOP代理，方法入口就开启事务',
    '订单创建 + 商品状态更新 + 延时消息发送在同一事务',
    'RocketMQ delayTimeLevel=16 \u2192 30分钟超时',
], sz=11, bullet_color=C3)

rect(s, 7.0, 4.6, 5.5, 2.6, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 7.2, 4.75, 0.05, 2.3)
txt(s, 7.5, 4.65, 5.0, 0.35, '支付安全三步验证', sz=15, color=C3, bold=True)
bullets(s, 7.5, 5.05, 5.0, 2.0, [
    '\u2460 RSA验签 \u2014 AlipaySignature.rsaCheckV1()',
    '\u2461 幂等性检查 \u2014 检查paymentStatus，已支付直接返回',
    '\u2462 双重回调保障 \u2014 同步return + 异步notify',
    '本地开发无公网IP，异步回调可能无法到达',
], sz=11, bullet_color=C3)
page_num(s, 15, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 16: KEY TECH - ES SEARCH
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '04', '关键技术 \u2014 Elasticsearch搜索', 'BoolQuery多层嵌套 + IK分词 + 高亮')
rect(s, 0.8, 1.5, 5.8, 5.5, C2, radius=5000)
txt(s, 1.1, 1.7, 5.3, 0.35, '核心查询代码', sz=16, color=C3, bold=True)
code_text = (
    '// ItemServiceImpl.search()\n'
    'Query q = NativeQuery.builder()\n'
    '  .withQuery(q -> q.bool(b -> {\n'
    '    // \u5fc5\u987b\u5728\u552e\n'
    '    b.must(m -> m.term(t ->\n'
    '      t.field("status")\n'
    '       .value(1)));\n'
    '    // \u5173\u952e\u8bcd\u641c\u7d22\n'
    '    b.must(m -> m.bool(kb -> {\n'
    '      kb.should(s -> s.match(\n'
    '        ma -> ma.field("title")\n'
    '          .query(keyword)));\n'
    '      kb.should(s -> s.match(\n'
    '        ma -> ma.field("description")\n'
    '          .query(keyword)));\n'
    '    }));\n'
    '  }))\n'
    '  .withHighlightQuery(\n'
    '    new HighlightQuery(...))\n'
    '  .build();'
)
txt(s, 1.1, 2.2, 5.3, 4.5, code_text, sz=11, color=C5, font='Consolas')

rect(s, 7.0, 1.5, 5.5, 2.5, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 7.2, 1.65, 0.05, 2.2)
txt(s, 7.5, 1.55, 5.0, 0.35, '搜索功能特性', sz=16, color=C3, bold=True)
bullets(s, 7.5, 1.95, 5.0, 2.0, [
    'BoolQuery 多层嵌套 DSL',
    'IK 分词 + 高亮渲染（HighlightQuery）',
    '支持通配符查询（wildcard）',
    '四种排序：最新/价格升/价格降/最热',
], sz=12, bullet_color=C3)

rect(s, 7.0, 4.2, 5.5, 2.8, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 7.2, 4.35, 0.05, 2.5)
txt(s, 7.5, 4.25, 5.0, 0.35, '自动补全方案变更', sz=16, color=ORANGE, bold=True)
bullets(s, 7.5, 4.65, 5.0, 2.2, [
    '最初：Completion Suggester',
    '问题：要求字段类型为 completion，与 text 冲突',
    '重新创建索引影响线上服务',
    '降级：match_phrase_prefix 查询 title 字段',
    '稳定且无需额外映射配置',
    '返回前10条匹配结果，下拉列表展示',
], sz=12, bullet_color=ORANGE)
page_num(s, 16, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 17: TESTING - FUNCTIONAL
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '05', '系统测试 \u2014 功能测试', '6项核心测试用例全部通过')
tests = [
    ('UC-01','登录与注册','普通用户/管理员登录\n空密码/错误密码异常提示','\u2713'),
    ('UC-02','商品智能发布','AI润色合格率95%+\n发布流程完整','\u2713'),
    ('UC-03','AI导购与检索','RAG检索真实商品\n推荐可点击跳转','\u2713'),
    ('UC-04','线上交易订单','下单\u2192支付\u2192确认收货\n完整链路状态迁移正确','\u2713'),
    ('UC-05','订单超时取消','RocketMQ延时消息\n30分钟自动取消+恢复在售','\u2713'),
    ('UC-06','AES加密与越权防护','数据库密文/API明文\n/admin/**返回403','\u2713'),
]
for i,(uc,name,desc,result) in enumerate(tests):
    r = i // 2; c = i % 2
    x = 0.8 + c*6.2; y = 1.5 + r*1.85
    rect(s, x, y, 5.8, 1.65, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
    rect(s, x, y, 1.0, 1.65, C2, radius=5000)
    txt(s, x+0.05, y+0.15, 0.9, 0.3, uc, sz=10, color=C4, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x+0.05, y+0.5, 0.9, 0.3, name, sz=10, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, x+1.2, y+0.15, 3.5, 1.3, desc, sz=11, color=GRAY)
    rect(s, x+4.8, y+0.3, 0.8, 0.8, C6, radius=5000)
    txt(s, x+4.8, y+0.35, 0.8, 0.7, result, sz=24, color=C3, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
page_num(s, 17, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 18: TESTING - PERFORMANCE
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '05', '系统测试 \u2014 性能测试', 'JMeter 200并发线程持续施压5分钟')
rect(s, 0.8, 1.5, 7.5, 0.5, C2, radius=3000)
txt(s, 1.0, 1.55, 2.5, 0.4, '接口', sz=13, color=WHITE, bold=True)
txt(s, 3.5, 1.55, 1.5, 0.4, '平均响应', sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 5.0, 1.55, 1.5, 0.4, 'TPS', sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 6.5, 1.55, 1.5, 0.4, '99%延迟', sz=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
perf = [
    ('首页商品列表(ES)','45ms','>2100','<100ms',C3),
    ('商品搜索+高亮','60ms','~1800','<150ms',C4),
    ('AI导购接口','3.5s','~15','~5s',ORANGE),
]
for i,(name,resp,tps,p99,clr) in enumerate(perf):
    y = 2.1 + i*0.6
    bg_c = WHITE if i%2==0 else C7
    rect(s, 0.8, y, 7.5, 0.5, bg_c)
    txt(s, 1.0, y+0.05, 2.5, 0.4, name, sz=12, color=DARK, bold=True)
    txt(s, 3.5, y+0.05, 1.5, 0.4, resp, sz=14, color=clr, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 5.0, y+0.05, 1.5, 0.4, tps, sz=14, color=clr, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 6.5, y+0.05, 1.5, 0.4, p99, sz=12, color=GRAY, align=PP_ALIGN.CENTER)

rect(s, 0.8, 4.0, 7.5, 1.5, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 4.15, 0.05, 1.2)
txt(s, 1.3, 4.1, 6.8, 0.35, '性能结论', sz=16, color=C3, bold=True)
bullets(s, 1.3, 4.5, 6.8, 1.0, [
    '常规读写接口在200并发压力下表现稳定，响应时间和吞吐量均满足设计目标',
    'AI接口瓶颈在模型推理环节（~3s），后续可通过流式输出和模型蒸馏优化',
    '测试过程中未发现内存泄漏或线程死锁问题',
], sz=11, bullet_color=C3)

txt(s, 9.0, 1.5, 3.8, 0.4, '测试环境', sz=18, color=C3, bold=True)
rect(s, 9.0, 2.0, 3.8, 4.5, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 9.2, 2.15, 0.04, 4.2)
env_items = [
    '4核 8GB CentOS 7',
    'JDK 17 + MySQL 8.0',
    'Redis 7.0 + ES 8.12',
    'PostgreSQL 16 + PgVector',
    'RocketMQ 5.1',
    'Nginx 反向代理',
    'JMeter 压测工具',
]
bullets(s, 9.4, 2.15, 3.2, 4.2, env_items, sz=11, bullet_color=C3)
page_num(s, 18, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 19: TESTING - SECURITY
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '05', '系统测试 \u2014 安全测试', 'AES加密 + 越权防护')
img(s, 'img_999c829c.png', 0.3, 1.4, h=2.5)
img(s, 'img_5b095e3f.png', 4.5, 1.4, h=2.5)
img(s, 'img_506da320.png', 8.7, 1.4, h=2.5)

rect(s, 0.8, 4.2, 5.8, 2.8, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 1.0, 4.35, 0.05, 2.5)
txt(s, 1.3, 4.25, 5.0, 0.35, 'AES字段级加密验证', sz=16, color=C3, bold=True)
bullets(s, 1.3, 4.65, 5.0, 2.2, [
    '直接查询数据库：phone字段为Base64密文',
    '通过API接口查询：phone返回明文',
    'TypeHandler加解密逻辑正确',
    '即使数据库泄露也无法获取手机号明文',
], sz=12, bullet_color=C3)

rect(s, 7.0, 4.2, 5.5, 2.8, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
accent_bar(s, 7.2, 4.35, 0.05, 2.5)
txt(s, 7.5, 4.25, 5.0, 0.35, '越权防护验证', sz=16, color=C3, bold=True)
bullets(s, 7.5, 4.65, 5.0, 2.2, [
    '普通用户访问 /admin/** \u2192 返回403',
    '普通用户删除他人商品 \u2192 "无权操作"',
    'AdminInterceptor 校验 role=1',
    '权限校验机制工作正常',
], sz=12, bullet_color=C3)
page_num(s, 19, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 20: SUMMARY
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '06', '总结与展望', '主要工作与创新点')
txt(s, 0.8, 1.5, 5.5, 0.4, '主要工作', sz=20, color=C3, bold=True)
works = [
    ('基础交易链路','Spring Boot + Vue 3 完整闭环\n发布 \u2192 搜索 \u2192 通讯 \u2192 交易 \u2192 支付'),
    ('AI深度整合','发布侧AI润色降低门槛（合格率95%+）\n导购侧RAG管线提升匹配效率'),
    ('工程落地','双轨同步保证向量数据一致性\nPrompt迭代调优收敛AI输出行为'),
]
for i,(title,desc) in enumerate(works):
    y = 2.0 + i*1.6
    rect(s, 0.8, y, 5.5, 1.4, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
    rect(s, 0.8, y, 0.7, 1.4, C3, radius=5000)
    txt(s, 0.85, y+0.2, 0.6, 0.5, str(i+1), sz=24, color=WHITE, bold=True, font='Arial Black', align=PP_ALIGN.CENTER)
    txt(s, 1.7, y+0.1, 4.4, 0.35, title, sz=15, color=C3, bold=True)
    txt(s, 1.7, y+0.5, 4.4, 0.8, desc, sz=12, color=GRAY)

txt(s, 7.0, 1.5, 5.5, 0.4, '改进方向', sz=20, color=C3, bold=True)
improvements = [
    ('推理延迟优化','同步 \u2192 流式输出 chatClient.stream()\n向量检索与Prompt拼接并行化\n部署轻量蒸馏模型'),
    ('多模态检索','引入CLIP视觉-语言模型\n商品图片生成Embedding\n实现以图搜图能力'),
    ('移动端适配','响应式布局 / 微信小程序\n覆盖校园手机使用场景\n考虑开发独立移动端页面'),
]
for i,(title,desc) in enumerate(improvements):
    y = 2.0 + i*1.6
    rect(s, 7.0, y, 5.5, 1.4, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
    rect(s, 7.0, y, 0.7, 1.4, ORANGE, radius=5000)
    txt(s, 7.05, y+0.2, 0.6, 0.5, str(i+1), sz=24, color=WHITE, bold=True, font='Arial Black', align=PP_ALIGN.CENTER)
    txt(s, 7.9, y+0.1, 4.4, 0.35, title, sz=15, color=ORANGE, bold=True)
    txt(s, 7.9, y+0.5, 4.4, 0.8, desc, sz=12, color=GRAY)
page_num(s, 20, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 21: Q&A
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C7)
section_header(s, '', '高频答辩问题准备', '15道核心Q&A')
qa_pairs = [
    ('Q1','为什么选JWT而不是Session？','JWT无状态，水平扩展无需共享Session'),
    ('Q2','RAG如何防止AI幻觉？','Prompt硬性指令 + RAG事实锚点'),
    ('Q3','向量数据如何与MySQL保持一致？','双轨同步：增量保实时，全量保正确'),
    ('Q4','为什么用TransactionTemplate？','业务校验后才开启事务，避免长事务'),
    ('Q5','similarityThreshold为什么设0.0？','校园商品库有限，避免查询无果'),
    ('Q6','为什么不用Spring Security？','权限模型简单，自定义拦截器更轻量'),
    ('Q7','AES为什么用TypeHandler？','对业务代码零侵入，注解即可'),
    ('Q8','自动补全为什么用match_phrase_prefix？','Completion Suggester mapping冲突'),
]
for i,(q,title,ans) in enumerate(qa_pairs):
    r = i // 4; c = i % 4
    x = 0.6 + c*3.15; y = 1.5 + r*1.85
    rect(s, x, y, 2.95, 1.65, WHITE, line=RGBColor(0xE2,0xE8,0xF0), radius=5000)
    rect(s, x, y, 2.95, 0.4, C2, radius=5000)
    txt(s, x+0.1, y+0.05, 0.5, 0.3, q, sz=10, color=C3, bold=True)
    txt(s, x+0.6, y+0.05, 2.2, 0.3, title, sz=10, color=WHITE, bold=True)
    txt(s, x+0.15, y+0.5, 2.6, 1.0, ans, sz=10, color=GRAY)
page_num(s, 21, TOTAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 22: THANK YOU
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, C1)
rect(s, 0, 0, 13.333, 0.06, C3)
rect(s, 0, 7.44, 13.333, 0.06, C3)
txt(s, 1.5, 2.0, 10, 0.8, '感谢各位老师', sz=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1.5, 3.0, 10, 0.5, '敬请批评指正', sz=24, color=C4, align=PP_ALIGN.CENTER)
rect(s, 5.5, 3.8, 2.3, 0.02, C3)
txt(s, 1.5, 4.2, 10, 0.4, '答辩人：刘遥杰', sz=16, color=C4, align=PP_ALIGN.CENTER)
txt(s, 1.5, 4.7, 10, 0.4, '基于RAG的校园智能二手交易系统设计与实现', sz=14, color=LGRAY, align=PP_ALIGN.CENTER)
txt(s, 1.5, 5.3, 10, 0.3, '武汉理工大学  计算机与人工智能学院', sz=13, color=LGRAY, align=PP_ALIGN.CENTER)

out = r'd:\SwapU\docs\thesis\SwapU答辩PPT.pptx'
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
