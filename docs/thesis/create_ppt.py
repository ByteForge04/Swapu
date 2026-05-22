import sys
sys.path.insert(0, r'D:\SwapU\docs\thesis\pip_libs')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT = RGBColor(0x08, 0x91, 0xB2)
ACCENT2 = RGBColor(0x06, 0x5A, 0x82)

IMG_DIR = r'd:\SwapU\docs\thesis\thesis_images'

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, x, y, w, h, text, font_size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_list(slide, x, y, w, h, items, font_size=16, color=DARK, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_image_safe(slide, img_name, x, y, w=None, h=None):
    path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(path):
        kwargs = {}
        if w: kwargs['width'] = Inches(w)
        if h: kwargs['height'] = Inches(h)
        slide.shapes.add_picture(path, Inches(x), Inches(y), **kwargs)

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 1.5, 1.8, 10, 1.2, '基于RAG的校园智能二手交易系统', font_size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name='Calibri')
add_text(slide, 1.5, 3.0, 10, 0.6, '设计与实现', font_size=28, color=ICE, bold=False, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.2, 10, 0.5, 'SwapU — 让校园闲置流转更智能', font_size=18, color=ICE, align=PP_ALIGN.CENTER)
add_rect(slide, 5.5, 5.0, 2.3, 0.03, ACCENT)
add_text(slide, 1.5, 5.3, 10, 0.4, '答辩人：刘遥杰    指导教师：霞', font_size=16, color=ICE, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 5.8, 10, 0.4, '计算机与人工智能学院  计算机科学与技术专业 zy2201班', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Outline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 5, 0.6, '目录', font_size=32, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.1, 1.5, 0.04, ACCENT)

sections = [
    ('01', '选题背景与意义', '校园闲置错配 + 现有平台痛点'),
    ('02', '系统架构设计', '五层架构 + 四种存储引擎'),
    ('03', '核心功能展示', 'AI润色 / RAG导购 / 订单交易'),
    ('04', '关键技术实现', 'JWT / AES / 双轨同步 / 状态机'),
    ('05', '系统测试', '功能 / 性能 / 安全三维度验证'),
    ('06', '总结与展望', '主要工作 + 改进方向'),
]
for i, (num, title, desc) in enumerate(sections):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.1
    y = 1.8 + row * 2.5
    add_rect(slide, x, y, 3.8, 2.0, LIGHT_BG)
    add_rect(slide, x, y, 0.08, 2.0, ACCENT)
    add_text(slide, x + 0.3, y + 0.2, 3.3, 0.5, num, font_size=28, color=ACCENT, bold=True)
    add_text(slide, x + 0.3, y + 0.8, 3.3, 0.5, title, font_size=18, color=DARK, bold=True)
    add_text(slide, x + 0.3, y + 1.3, 3.3, 0.5, desc, font_size=12, color=GRAY)

# ============================================================
# SLIDE 3: Background
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '01  选题背景与意义', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_text(slide, 0.8, 1.3, 5.5, 0.4, '现实痛点', font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, 0.8, 1.8, 5.5, 1.5, [
    '• 每学期末大量教材、电子设备被淘汰',
    '• 新生群体年复一年存在刚性采购需求',
    '• 二者之间存在结构性错配',
], font_size=14)

add_text(slide, 0.8, 3.3, 5.5, 0.4, '现有平台三方面错位', font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, 0.8, 3.8, 5.5, 2.0, [
    '• 信息匹配效率不足 — 学生关心相邻宿舍楼而非全国信用分',
    '• 交易摩擦过大 — 校园以当面交易为主，平台缺乏面交支持',
    '• 发布门槛偏高 — 学生不愿为短期闲置品撰写详细描述',
], font_size=14)

add_text(slide, 7.0, 1.3, 5.5, 0.4, 'SwapU 的解决思路', font_size=20, color=ACCENT, bold=True)
add_rect(slide, 7.0, 1.8, 5.5, 4.5, LIGHT_BG)
add_rect(slide, 7.0, 1.8, 0.08, 4.5, ACCENT)
add_bullet_list(slide, 7.3, 2.0, 5.0, 4.0, [
    '发布侧 → AI 一键润色，降低文案撰写门槛',
    '导购侧 → RAG 管线，理解口语化购物意图',
    '交易侧 → 校园当面交易，降低交易摩擦',
    '',
    '核心创新：将大语言模型深度整合进业务流程',
    '而非仅仅作为噱头简单接入',
], font_size=14)

# ============================================================
# SLIDE 4: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '02  系统架构设计', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_image_safe(slide, 'img_ce8dc44d.png', 0.5, 1.3, w=6.5)

add_text(slide, 7.5, 1.3, 5.3, 0.4, '五层架构', font_size=20, color=ACCENT, bold=True)
layers = [
    ('表现层', 'Vue 3 + Element Plus + Pinia'),
    ('网关控制层', 'LoginInterceptor + AdminInterceptor'),
    ('业务逻辑层', '15 Controller + 13 ServiceImpl'),
    ('数据持久层', 'MySQL + Redis + ES + PgVector'),
    ('基础设施层', 'RocketMQ + 支付宝 + Ollama + DeepSeek'),
]
for i, (name, desc) in enumerate(layers):
    y = 1.9 + i * 0.9
    add_rect(slide, 7.5, y, 5.3, 0.75, LIGHT_BG)
    add_rect(slide, 7.5, y, 0.08, 0.75, ACCENT)
    add_text(slide, 7.8, y + 0.05, 2.0, 0.35, name, font_size=14, color=ACCENT, bold=True)
    add_text(slide, 7.8, y + 0.38, 4.8, 0.35, desc, font_size=12, color=GRAY)

add_text(slide, 7.5, 6.5, 5.3, 0.4, '各存储引擎职责隔离，互不重叠', font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Feature - Login & Home
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '03  核心功能展示 — 登录与首页', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_image_safe(slide, 'img_f93bbfac.png', 0.5, 1.3, h=4.5)
add_image_safe(slide, 'img_a1b8e6be.png', 6.8, 1.3, h=4.5)

add_text(slide, 0.5, 6.0, 5.8, 0.5, '登录页面：BCrypt加密 + JWT Token签发', font_size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 6.8, 6.0, 5.8, 0.5, '首页：卡片式信息流 + 分类导航', font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: Feature - AI Polish & RAG Guide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '03  核心功能展示 — AI 智能交互', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_image_safe(slide, 'img_6f39e11f.png', 0.3, 1.3, h=4.5)
add_image_safe(slide, 'img_74beebc6.png', 6.5, 1.3, h=4.5)

add_text(slide, 0.3, 6.0, 6.0, 0.5, 'AI润色：输入关键词 → DeepSeek生成结构化文案', font_size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 6.5, 6.0, 6.5, 0.5, 'AI导购：自然语言 → RAG检索 → 推荐真实在售商品', font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: Feature - Chat & Orders
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '03  核心功能展示 — 通讯与订单', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_image_safe(slide, 'img_ae0637aa.png', 0.3, 1.3, h=4.5)
add_image_safe(slide, 'img_a0a7b452.png', 6.5, 1.3, h=4.5)

add_text(slide, 0.3, 6.0, 6.0, 0.5, 'WebSocket即时通讯：文本+图片消息', font_size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 6.5, 6.0, 6.5, 0.5, '订单管理：买入/卖出双视角 + 状态流转', font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: Feature - Admin
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '03  核心功能展示 — 后台管理', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_image_safe(slide, 'img_ecb8add6.png', 2.0, 1.3, h=4.8)

add_text(slide, 2.0, 6.3, 9.0, 0.5, '管理员Dashboard：ECharts运营指标 + 商品审核 + 举报处理 + 公告管理', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: Key Tech - JWT & AES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '04  关键技术 — JWT认证 & AES加密', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_rect(slide, 0.8, 1.4, 5.8, 5.3, LIGHT_BG)
add_rect(slide, 0.8, 1.4, 0.08, 5.3, ACCENT)
add_text(slide, 1.1, 1.5, 5.3, 0.4, 'JWT 无状态认证', font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, 1.1, 2.1, 5.3, 4.0, [
    '选型理由：无状态，水平扩展无需共享Session',
    '签发：Keys.hmacShaKeyFor() → HMAC-SHA256',
    'Claims：userId + username + role，24h有效期',
    '传输：Authorization 请求头（无 Bearer 前缀）',
    '验证：LoginInterceptor → validateToken()',
    '兼容：Integer/Long 类型判断（JVM版本差异）',
], font_size=13)

add_rect(slide, 6.9, 1.4, 5.8, 5.3, LIGHT_BG)
add_rect(slide, 6.9, 1.4, 0.08, 5.3, ACCENT)
add_text(slide, 7.2, 1.5, 5.3, 0.4, 'AES 字段加密', font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, 7.2, 2.1, 5.3, 4.0, [
    '方案：MyBatis-Plus TypeHandler 机制',
    '写入：AesEncryptUtils.encrypt() → Base64密文',
    '读取：AesEncryptUtils.decrypt() → 明文',
    '注解：@TableField(typeHandler = ...)',
    '效果：数据库密文 / API明文，业务零感知',
    '安全：即使数据库泄露也无法获取手机号',
], font_size=13)

# ============================================================
# SLIDE 10: Key Tech - RAG
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '04  关键技术 — RAG智能导购', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_image_safe(slide, 'img_a1392b1a.png', 0.3, 1.2, h=3.5)

add_text(slide, 7.0, 1.3, 5.8, 0.4, '三步链路', font_size=20, color=ACCENT, bold=True)

steps = [
    ('1', '向量检索', 'bge-m3 生成1024维向量\nPgVector 余弦相似度 topK=3'),
    ('2', '上下文拼接', '检索结果注入 System Prompt\n历史对话截取最近6条'),
    ('3', '推理与封装', 'chatClient.call() 同步推理\n返回 AI文本 + 商品卡片'),
]
for i, (num, title, desc) in enumerate(steps):
    y = 1.9 + i * 1.5
    add_rect(slide, 7.0, y, 5.8, 1.3, LIGHT_BG)
    add_rect(slide, 7.0, y, 0.08, 1.3, ACCENT)
    add_text(slide, 7.3, y + 0.05, 0.5, 0.4, num, font_size=24, color=ACCENT, bold=True)
    add_text(slide, 7.9, y + 0.05, 2.0, 0.4, title, font_size=16, color=DARK, bold=True)
    add_text(slide, 7.9, y + 0.45, 4.7, 0.8, desc, font_size=12, color=GRAY)

add_text(slide, 7.0, 6.4, 5.8, 0.5, '防幻觉：Prompt硬性指令 + RAG事实锚点', font_size=13, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11: Key Tech - Dual Sync
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '04  关键技术 — RAG双轨同步策略', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_image_safe(slide, 'img_dd555ab5.png', 0.3, 1.2, h=3.8)

add_rect(slide, 7.0, 1.3, 5.8, 2.3, LIGHT_BG)
add_rect(slide, 7.0, 1.3, 0.08, 2.3, ACCENT)
add_text(slide, 7.3, 1.4, 5.3, 0.4, '第一轨：增量同步', font_size=16, color=ACCENT, bold=True)
add_bullet_list(slide, 7.3, 1.9, 5.3, 1.5, [
    'RocketMQ 消息驱动',
    'MySQL事务完成 → 发送ItemSyncMessage',
    'ItemSyncListener 消费 → ES + PgVector同步',
    'MQ不可用时降级为同步直连',
], font_size=12)

add_rect(slide, 7.0, 3.8, 5.8, 2.3, LIGHT_BG)
add_rect(slide, 7.0, 3.8, 0.08, 2.3, ACCENT)
add_text(slide, 7.3, 3.9, 5.3, 0.4, '第二轨：全量同步', font_size=16, color=ACCENT, bold=True)
add_bullet_list(slide, 7.3, 4.4, 5.3, 1.5, [
    '管理员触发 /admin/es/sync',
    '遍历所有 status=1 的商品',
    '批量写入 ES 和 PgVector',
    '用于系统初始化或灾难恢复',
], font_size=12)

add_text(slide, 7.0, 6.3, 5.8, 0.5, '核心思想：增量保实时，全量保正确', font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: Key Tech - Order State Machine
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '04  关键技术 — 订单状态机与支付', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_image_safe(slide, 'img_8370a17d.png', 0.3, 1.2, h=3.2)
add_image_safe(slide, 'img_2514cc9f.png', 6.5, 1.2, h=3.2)

add_rect(slide, 0.8, 4.7, 5.5, 2.3, LIGHT_BG)
add_rect(slide, 0.8, 4.7, 0.08, 2.3, ACCENT)
add_text(slide, 1.1, 4.8, 5.0, 0.4, '编程式事务', font_size=16, color=ACCENT, bold=True)
add_bullet_list(slide, 1.1, 5.3, 5.0, 1.5, [
    'TransactionTemplate 而非 @Transactional',
    '业务校验后才开启事务，避免长事务',
    'RocketMQ 延时消息 30分钟超时自动取消',
], font_size=12)

add_rect(slide, 6.9, 4.7, 5.8, 2.3, LIGHT_BG)
add_rect(slide, 6.9, 4.7, 0.08, 2.3, ACCENT)
add_text(slide, 7.2, 4.8, 5.3, 0.4, '支付安全三步验证', font_size=16, color=ACCENT, bold=True)
add_bullet_list(slide, 7.2, 5.3, 5.3, 1.5, [
    '① RSA验签 — AlipaySignature.rsaCheckV1()',
    '② 幂等性检查 — 检查paymentStatus',
    '③ 双重回调 — 同步return + 异步notify',
], font_size=12)

# ============================================================
# SLIDE 13: Testing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '05  系统测试', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_text(slide, 0.8, 1.3, 4.0, 0.4, '功能测试', font_size=20, color=ACCENT, bold=True)
func_tests = [
    ('登录与注册', '✅ 通过'),
    ('商品智能发布', '✅ 通过'),
    ('AI导购与检索', '✅ 通过'),
    ('线上交易订单', '✅ 通过'),
    ('订单超时取消', '✅ 通过'),
    ('AES加密与越权防护', '✅ 通过'),
]
for i, (name, result) in enumerate(func_tests):
    y = 1.8 + i * 0.55
    add_rect(slide, 0.8, y, 5.5, 0.45, LIGHT_BG)
    add_text(slide, 1.0, y + 0.05, 3.5, 0.35, name, font_size=13, color=DARK)
    add_text(slide, 4.5, y + 0.05, 1.5, 0.35, result, font_size=13, color=RGBColor(0x16, 0xA3, 0x4A), bold=True, align=PP_ALIGN.RIGHT)

add_text(slide, 7.0, 1.3, 5.5, 0.4, '性能测试（200并发）', font_size=20, color=ACCENT, bold=True)
perf_data = [
    ('首页商品列表', '45ms', '>2100'),
    ('商品搜索', '60ms', '~1800'),
    ('AI导购接口', '3.5s', '~15'),
]
add_rect(slide, 7.0, 1.8, 5.5, 0.4, RGBColor(0x1E, 0x27, 0x61))
add_text(slide, 7.1, 1.85, 2.0, 0.3, '接口', font_size=12, color=WHITE, bold=True)
add_text(slide, 9.0, 1.85, 1.5, 0.3, '平均响应', font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 10.5, 1.85, 1.8, 0.3, 'TPS', font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, (name, resp, tps) in enumerate(perf_data):
    y = 2.3 + i * 0.45
    bg_c = LIGHT_BG if i % 2 == 0 else WHITE
    add_rect(slide, 7.0, y, 5.5, 0.4, bg_c)
    add_text(slide, 7.1, y + 0.05, 2.0, 0.3, name, font_size=12, color=DARK)
    add_text(slide, 9.0, y + 0.05, 1.5, 0.3, resp, font_size=12, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, 10.5, y + 0.05, 1.8, 0.3, tps, font_size=12, color=DARK, align=PP_ALIGN.CENTER)

add_text(slide, 7.0, 3.8, 5.5, 0.4, '安全测试', font_size=20, color=ACCENT, bold=True)
add_rect(slide, 7.0, 4.3, 5.5, 1.5, LIGHT_BG)
add_rect(slide, 7.0, 4.3, 0.08, 1.5, ACCENT)
add_bullet_list(slide, 7.3, 4.4, 5.0, 1.3, [
    '✅ AES加密：数据库密文 / API明文',
    '✅ 越权防护：/admin/** 返回403',
    '✅ 删除他人商品返回"无权操作"',
], font_size=13)

add_text(slide, 0.8, 5.5, 5.5, 0.4, '已知缺陷', font_size=16, color=RGBColor(0xDC, 0x26, 0x26), bold=True)
add_bullet_list(slide, 0.8, 6.0, 5.5, 1.0, [
    '登录接口未配置防护策略，存在暴力破解风险',
    '建议：每分钟5次失败后锁定账号15分钟',
], font_size=12, color=GRAY)

# ============================================================
# SLIDE 14: Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 0.8, 0.4, 8, 0.6, '06  总结与展望', font_size=28, color=NAVY, bold=True)
add_rect(slide, 0.8, 1.0, 1.5, 0.04, ACCENT)

add_text(slide, 0.8, 1.3, 5.5, 0.4, '主要工作', font_size=20, color=ACCENT, bold=True)
works = [
    ('基础交易链路', 'Spring Boot + Vue 3 完整闭环\n发布→搜索→通讯→交易→支付'),
    ('AI深度整合', '发布侧AI润色降低门槛\n导购侧RAG管线提升匹配效率'),
    ('工程落地', '双轨同步保证向量数据一致性\nPrompt迭代调优收敛AI输出行为'),
]
for i, (title, desc) in enumerate(works):
    y = 1.8 + i * 1.5
    add_rect(slide, 0.8, y, 5.5, 1.3, LIGHT_BG)
    add_rect(slide, 0.8, y, 0.08, 1.3, ACCENT)
    add_text(slide, 1.1, y + 0.05, 5.0, 0.4, title, font_size=16, color=ACCENT, bold=True)
    add_text(slide, 1.1, y + 0.5, 5.0, 0.7, desc, font_size=12, color=GRAY)

add_text(slide, 7.0, 1.3, 5.5, 0.4, '改进方向', font_size=20, color=ACCENT, bold=True)
improvements = [
    ('推理延迟优化', '同步→流式输出 chatClient.stream()\n向量检索与Prompt拼接并行化'),
    ('多模态检索', '引入CLIP视觉-语言模型\n实现以图搜图能力'),
    ('移动端适配', '响应式布局 / 微信小程序\n覆盖校园手机使用场景'),
]
for i, (title, desc) in enumerate(improvements):
    y = 1.8 + i * 1.5
    add_rect(slide, 7.0, y, 5.5, 1.3, LIGHT_BG)
    add_rect(slide, 7.0, y, 0.08, 1.3, ACCENT)
    add_text(slide, 7.3, y + 0.05, 5.0, 0.4, title, font_size=16, color=ACCENT, bold=True)
    add_text(slide, 7.3, y + 0.5, 5.0, 0.7, desc, font_size=12, color=GRAY)

# ============================================================
# SLIDE 15: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
add_text(slide, 1.5, 2.2, 10, 1.0, '感谢各位老师', font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 3.4, 10, 0.6, '敬请批评指正', font_size=24, color=ICE, align=PP_ALIGN.CENTER)
add_rect(slide, 5.5, 4.3, 2.3, 0.03, ACCENT)
add_text(slide, 1.5, 4.8, 10, 0.4, '答辩人：刘遥杰', font_size=16, color=ICE, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 5.3, 10, 0.4, '基于RAG的校园智能二手交易系统设计与实现', font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

output_path = r'd:\SwapU\docs\thesis\SwapU答辩PPT.pptx'
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
