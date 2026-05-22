#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Generate defense PPT for SwapU thesis.
Wuhan University of Technology - Undergraduate Thesis Defense
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== Config =====
OUTPUT_PATH = r'd:\SwapU\docs\thesis\SwapU答辩PPT.pptx'
TITLE_FONT = '微软雅黑'
BODY_FONT = '微软雅黑'
TITLE_SIZE = Pt(32)
SUBTITLE_SIZE = Pt(20)
BODY_SIZE = Pt(16)
SMALL_SIZE = Pt(12)

# Color scheme
BLUE_DARK = RGBColor(0x1B, 0x3A, 0x5C)   # Deep navy
BLUE_MID = RGBColor(0x2C, 0x5F, 0x8A)    # Mid blue
BLUE_LIGHT = RGBColor(0x3A, 0x7C, 0xBF)  # Light blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x88, 0x88, 0x88)
ACCENT = RGBColor(0xE8, 0x6A, 0x17)      # Orange accent

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ===== Helper Functions =====

def add_bg(slide, color=BLUE_DARK):
    """Add solid color background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=BODY_SIZE,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name=BODY_FONT):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_title_slide(prs, title, subtitle=""):
    """Create a title slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, BLUE_DARK)

    # Accent line
    add_shape_bg(slide, Inches(1), Inches(2.8), Inches(1.5), Pt(4), ACCENT)

    # Title
    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 title, Pt(40), WHITE, True, PP_ALIGN.LEFT)

    # Subtitle
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
                     subtitle, Pt(20), RGBColor(0xAA, 0xBB, 0xCC), False, PP_ALIGN.LEFT)

    # Footer
    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 '武汉理工大学 · 计算机与人工智能学院 · 2026届本科毕业设计答辩',
                 Pt(12), RGBColor(0x77, 0x88, 0x99), False, PP_ALIGN.LEFT)
    return slide

def add_content_slide(prs, title, content_lines, emphasis_lines=None):
    """Create a content slide with left title bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, WHITE)

    # Top bar
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(3), ACCENT)

    # Left accent bar
    add_shape_bg(slide, Inches(0.6), Inches(0.8), Pt(5), Inches(1.0), BLUE_DARK)

    # Title
    add_text_box(slide, Inches(0.9), Inches(0.7), Inches(11), Inches(0.8),
                 title, TITLE_SIZE, BLUE_DARK, True, PP_ALIGN.LEFT)

    # Separator line
    add_shape_bg(slide, Inches(0.9), Inches(1.5), Inches(11.5), Pt(1), RGBColor(0xDD, 0xDD, 0xDD))

    # Content
    y = Inches(1.8)
    for line in content_lines:
        if line.startswith('##'):
            # Sub-heading
            add_text_box(slide, Inches(0.9), y, Inches(11.5), Inches(0.5),
                         line.replace('##', '').strip(), Pt(20), BLUE_MID, True)
            y += Inches(0.5)
        elif line.startswith('- '):
            add_text_box(slide, Inches(1.3), y, Inches(11), Inches(0.4),
                         line[2:], BODY_SIZE, BLACK, False)
            y += Inches(0.4)
        elif line == '':
            y += Inches(0.15)
        else:
            add_text_box(slide, Inches(0.9), y, Inches(11.5), Inches(0.4),
                         line, BODY_SIZE, BLACK, False)
            y += Inches(0.4)

    # Emphasis box at the bottom
    if emphasis_lines:
        y += Inches(0.2)
        box = add_shape_bg(slide, Inches(0.9), y, Inches(11.5), Inches(len(emphasis_lines) * 0.35 + 0.4),
                           RGBColor(0xEE, 0xF3, 0xF9))
        y += Inches(0.15)
        for line in emphasis_lines:
            add_text_box(slide, Inches(1.1), y, Inches(11), Inches(0.35),
                         line, Pt(14), BLUE_DARK, True)
            y += Inches(0.35)

    # Page number
    add_text_box(slide, Inches(12), Inches(7.1), Inches(1), Inches(0.3),
                 str(prs.slides.index(slide) + 1), Pt(10), GRAY, False, PP_ALIGN.RIGHT)

    return slide

def add_code_slide(prs, title, code_lines, note=""):
    """Create a code-focused slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(0x1E, 0x1E, 0x1E))  # Dark code background

    # Top bar
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Pt(3), ACCENT)

    # Title in white
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
                 title, Pt(24), WHITE, True, PP_ALIGN.LEFT)

    # Code
    code_text = '\n'.join(code_lines)
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    p.font.name = 'Consolas'
    p.line_spacing = Pt(16)

    # Note/callout at bottom
    if note:
        add_text_box(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.8),
                     '💡 ' + note, Pt(16), RGBColor(0xFF, 0xCC, 0x66), True)

    # Page number
    add_text_box(slide, Inches(12), Inches(7.1), Inches(1), Inches(0.3),
                 str(prs.slides.index(slide) + 1), Pt(10), RGBColor(0x66, 0x66, 0x66), False, PP_ALIGN.RIGHT)

    return slide


# ===== SLIDE 1: 封面 =====
add_title_slide(prs,
    '基于 Vue3 + Spring Boot 的校园闲置物品\n循环利用系统设计与实现',
    'SwapU — 校园智能二手交易平台\n\n答辩人：XXX    指导教师：XXX')

# ===== SLIDE 2: 项目背景 =====
add_content_slide(prs, '项目背景与痛点', [
    '## 校园闲置物品的周期性积压',
    '- 每学期末：大量教材、电子产品、生活用品被闲置',
    '- 资源错配：一面是巨量闲置品产生，一面是新生刚性采购需求',
    '',
    '## 现有平台为什么没解决？',
    '- 信息匹配低效：全国大市场逻辑，无法按校区半径筛选',
    '- 发布门槛偏高：写吸引人的商品描述需要精力，平台无辅助',
    '- 沟通未场景化：当面交易为主的校园场景缺少本地化适配',
    '',
    '## SwapU 的设计思路',
    '- 用 AI 降低发布门槛 + 用 RAG 提升搜索精度 + 用 WebSocket 定制沟通',
], ['📌 核心命题：将大模型从"接入"变为"用好"，在垂直场景中工程化落地'])

# ===== SLIDE 3: 系统功能总览 =====
add_content_slide(prs, '系统功能总览', [
    '## 前台交易子系统',
    '- 用户认证：BCrypt加密注册登录 + JWT Token 鉴权',
    '- 商品发布：AI 文案润色 + AI 内容审核 → 一键生成结构化描述',
    '- 智能导购：自然语言提问 → PgVector 检索 → LLM 推荐真实商品',
    '- 全文搜索：Elasticsearch BoolQuery 复合查询 + 高亮 + 自动补全',
    '- 订单交易：Redisson 分布式锁防超卖 + 支付宝沙箱支付',
    '- 即时通讯：WebSocket 点对点实时聊天，消息持久化 MySQL',
    '',
    '## 后台管理子系统',
    '- 数据大盘：ECharts 近7日订单趋势 + 品类分布 + API调用频次',
    '- 内容审核：商品/留言/举报统一管理，审核操作闭环',
    '- 向量同步：RocketMQ 增量 + 全量补偿双轨策略',
], ['⚡ 15个Controller + 13个ServiceImpl + 1个WebSocket端点'])

# ===== SLIDE 4: 技术架构 =====
add_content_slide(prs, '系统技术架构（五层）', [
    '## 表现层',
    '- Nginx 反向代理：/api → 后端8080，静态资源自托管',
    '',
    '## 网关控制层',
    '- LoginInterceptor (JWT鉴权) → AdminInterceptor (管理员校验) → GlobalLogAspect (AOP审计)',
    '',
    '## 业务逻辑层',
    '- 15个Controller + 13个ServiceImpl + ChatWebSocketEndpoint',
    '- 核心Service：ItemServiceImpl（8个依赖注入）、TradeOrderServiceImpl（分布式锁+事务）',
    '',
    '## 数据持久层',
    '- MySQL 8.0（ACID写入）| Redis（缓存/锁/限流）| Elasticsearch 8.x（搜索副本）| PgVector（向量存储）',
    '',
    '## 基础设施层',
    '- RocketMQ（异步消息/延时任务）| 支付宝沙箱 | DeepSeek Chat | Ollama bge-m3',
], ['🔑 层间依赖严格控制，上层不碰下层数据实现——SOLID原则'])

# ===== SLIDE 5: RAG 智能导购 =====
add_content_slide(prs, '核心创新①：RAG 智能导购', [
    '## 核心流程："先检索，后生成"',
    '- ① 用户自然语言提问（如"两百以内有没有适合工科生的计算器"）',
    '- ② PgVector 余弦相似度检索（topK=3, threshold=0.0, HNSW索引）',
    '- ③ 将 Top 3 真实商品（标题、价格、描述）注入 System Prompt',
    '- ④ DeepSeek Chat 基于真实库存生成推荐 → 返回 { text, items }',
    '',
    '## 为什么这样设计？',
    '- 直接用LLM → "幻觉"：推荐平台上不存在的商品',
    '- RAG约束 → 模型只能推荐检索到的真实商品',
    '- topK=3 的实验依据：1→太窄，5→Prompt太长推荐变泛化',
    '- 截取最近6条历史（Math.max(0, history.size()-6)）保证多轮对话',
], ['🎯 核心设计决策：不是"把数据丢给模型"，而是"让模型回答被钉在真实库存上"'])

# ===== SLIDE 6: RAG 关键代码 =====
add_code_slide(prs, 'RAG 智能导购 — 关键代码', [
    '@PostMapping("/chat")',
    'public Result<Map<String, Object>> chat(@RequestBody Map request) {',
    '    // 第一步：RAG向量检索 topK=3',
    '    List<Map> items = ragService.retrieveRelevantItemDetails(msg, 3);',
    '',
    '    // 第二步：注入检索结果到System Prompt',
    '    systemPrompt.append("不要编造不存在的商品。\\n");',
    '    items.forEach(item -> systemPrompt.append(',
    '        "标题:" + item.title + " 价格:" + item.price + "\\n"));',
    '',
    '    // 第三步：截取最近6条历史 = 3轮对话',
    '    int start = Math.max(0, history.size() - 6);',
    '',
    '    // 第四步：同步调用LLM',
    '    String response = chatClient.call(prompt).getResult().getOutput().getContent();',
    '',
    '    // 第五步：封装返回 → AI文本 + 商品卡片数据',
    '    return Result.success(Map.of("text", response, "items", items));',
    '}',
], '关键：items 字段 = 检索到的真实商品 → 前端渲染可点击卡片 → 杜绝"幻觉推荐"')

# ===== SLIDE 7: 分布式锁 + 超时取消 =====
add_code_slide(prs, '核心创新②：分布式锁防超卖 + 延时消息超时取消', [
    'public Result<?> createOrder(TradeOrder order) {',
    '    // ① 分布式锁：商品粒度 (key = "item:lock:" + itemId)',
    '    RLock lock = redissonClient.getLock(lockKey);',
    '    lock.tryLock(3, 10, TimeUnit.SECONDS);  // 等待3s, 持有10s',
    '',
    '    // ② 编程式事务 (不用@Transactional!)',
    '    return transactionTemplate.execute(status -> {',
    '        // 校验status=1 → 保存订单status=0 → 更新商品status=2',
    '        item.setStatus(2); orderMapper.updateById(item);',
    '',
    '        // ③ 仅线上支付发送30分钟延时消息',
    '        if (order.getPaymentMethod() == 1) {',
    '            mqMsg.setDelayTimeLevel(16);  // level 16 = 30min',
    '            rocketMQTemplate.getProducer().send(mqMsg);',
    '        }',
    '        return Result.success(order.getOrderId());',
    '    });',
    '    // ④ finally { lock.unlock(); }  ← 精确控制释放时机',
    '}',
], '为什么TransactionTemplate？— 锁的释放必须在事务提交前，@Transactional无法控制')

# ===== SLIDE 8: 双轨数据同步 =====
add_content_slide(prs, '核心创新③：PgVector 与 MySQL 双轨同步', [
    '## 问题：向量库和关系库怎么保持一致？',
    '- 商品上架→向量没同步→用户搜不到',
    '- 商品售出→向量没删除→用户搜到已售出（比搜不到更糟）',
    '',
    '## 方案：增量同步 + 全量兜底',
    '',
    '## 第一轨：RocketMQ 增量同步',
    '- publish()中直接 CompletableFuture 异步写PgVector',
    '- 同时 syncToEs() → ItemSyncMessage → RocketMQ item-update-topic',
    '- → ItemSyncListener → doSyncToEs() → itemRepository.save(ES) + ragService.addDocument(PgVector)',
    '- 消费端重新查MySQL (getById) 而非用消息快照 — 避免异步延迟脏写',
    '',
    '## 第二轨：全量补偿',
    '- syncAllToEs() 手动触发，遍历所有status=1商品批量重建索引',
], ['🔄 两条写PgVector的路径存在冗余，但各自独立 — MQ延迟/单点故障时保证数据不丢'])

# ===== SLIDE 9: 数据库设计 =====
add_content_slide(prs, '核心数据库设计', [
    '## 关系型设计（MySQL 8.0, InnoDB）',
    '- sys_user：role=0用户/1管理员, phone字段AES-128/ECB/PKCS5Padding加密',
    '- item：五态状态机 (0待审核→1在售→2交易中→3已售出→4已下架)',
    '- trade_order：order_no=UUID去横线, amount=下单时快照价格',
    '- 关键索引：buyer_id, seller_id, user_id, category_id 均建B+树',
    '',
    '## 关键设计决策',
    '- amount 是下单时刻从 item.price 同步的快照 — 已写入不再随商品改价变动',
    '- payment_status 独立于 order_status — 线下交易可以不付款走完流程',
    '- Phone 通过 MyBatis-Plus TypeHandler 透明加解密 — Service层无感知',
    '',
    '## 向量存储（PostgreSQL + PgVector）',
    '- HNSW 索引 + COSINE_DISTANCE + 1024维 (bge-m3 输出)',
], ['📊 四核心表：sys_user → item → trade_order  (item_category)'])

# ===== SLIDE 10: 测试结果 =====
add_content_slide(prs, '系统测试结果', [
    '## 性能测试：500并发线程 × 5分钟',
    '- 首页商品列表：平均45ms, TPS 2100+',
    '- 商品搜索 (ES)：平均120ms, TPS ~850',
    '- AI 智能导购：平均3.5s (LLM推理2-3s + 网络 + Prompt处理)',
    '- CPU峰值 < 75%, 无OOM',
    '',
    '## 功能测试：全链路无断点',
    '- 发布→搜索→下单→支付→取消，7条路径状态迁移正确',
    '- RAG管线端到端验证：输入自然语言 → 检索 → LLM → 卡片渲染，无幻觉',
    '- 订单超时取消：30min后幂等校验通过，订单→status=3，商品→回滚status=1',
    '',
    '## 安全测试',
    '- AES加密：DB中phone为密文, API返回自动解密为明文 ✓',
    '- 限流：前3次200, 第4次起返回"下单过于频繁" ✓',
], ['⚠️ 已知不足：AI接口延迟3.5s偏高(待改SSE流式) | 登录接口未限流 | ES低配瓶颈'])

# ===== SLIDE 11: 总结与展望 =====
add_content_slide(prs, '总结与展望', [
    '## 主要工作',
    '- ① 完整交易链路：发布 → 搜索/AI导购 → 下单(锁+事务) → 支付 → 聊天',
    '- ② RAG 智能导购：三层约束(Prompt工程 + 事实约束 + 数据同步)抑制幻觉',
    '- ③ 工程可靠性：分布式锁防超卖 + 延时消息超时取消 + 双轨数据同步',
    '',
    '## 创新点回顾',
    '- 先检索后生成：不是把数据丢给模型，而是让模型被真实库存约束',
    '- 双轨同步：增量MQ + 全量兜底，工程级别的数据一致性方案',
    '- Prompt工程：十几轮迭代 → 三项硬约束 → 输出从"灾难"到"可靠"',
    '',
    '## 展望',
    '- ① SSE 流式输出：ChatController已import Flux，改.call()为.stream() → 首Token半秒内',
    '- ② 多模态检索：CLIP视觉模型 → "以图搜图" → 需GPU实例',
    '- ③ 补全安全短板：登录限流 + ES资源优化',
], ['🎓 达成设计目标：交易链路完整 + 中间件协同稳定 + AI能力工程化可用'])

# ===== SLIDE 12: 致谢 =====
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide12, BLUE_DARK)
add_shape_bg(slide12, Inches(5.5), Inches(3.0), Inches(2.5), Pt(4), ACCENT)
add_text_box(slide12, Inches(1), Inches(2.0), Inches(11.5), Inches(1.5),
             '感谢各位老师批评指正', Pt(44), WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide12, Inches(1), Inches(4.0), Inches(11.5), Inches(1),
             'SwapU — 校园智能二手交易平台', Pt(20), RGBColor(0xAA, 0xBB, 0xCC), False, PP_ALIGN.CENTER)
add_text_box(slide12, Inches(1), Inches(6.5), Inches(11.5), Inches(0.5),
             '武汉理工大学 · 计算机与人工智能学院 · 2026届本科毕业设计答辩',
             Pt(12), RGBColor(0x77, 0x88, 0x99), False, PP_ALIGN.CENTER)

# ===== Save =====
prs.save(OUTPUT_PATH)
print(f'PPT saved to: {OUTPUT_PATH}')
print(f'Total slides: {len(prs.slides)}')
